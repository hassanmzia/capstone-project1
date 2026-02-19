#!/usr/bin/env python3
"""Generate CNEAv5 Technical Architecture PowerPoint presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colors ──────────────────────────────────────────────────────────
BG_DARK      = RGBColor(0x0F, 0x11, 0x1A)
BG_SURFACE   = RGBColor(0x1A, 0x1D, 0x2E)
BG_CARD      = RGBColor(0x22, 0x26, 0x3A)
CYAN         = RGBColor(0x00, 0xD4, 0xFF)
PURPLE       = RGBColor(0xA855, 0xF7, 0x00)[0:3] if False else RGBColor(0xA8, 0x55, 0xF7)
GREEN        = RGBColor(0x22, 0xC5, 0x5E)
AMBER        = RGBColor(0xF5, 0x9E, 0x0B)
RED          = RGBColor(0xEF, 0x44, 0x44)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCC, 0xCC, 0xDD)
MID_GRAY     = RGBColor(0x88, 0x8A, 0xA0)
BORDER_COLOR = RGBColor(0x33, 0x37, 0x4D)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=12,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_arrow(slide, x1, y1, x2, y2, color=CYAN, width=Pt(2)):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = straight
    connector.line.color.rgb = color
    connector.line.width = width
    return connector


# ════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title Slide
# ════════════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_bg(slide1, BG_DARK)

# Title
add_text_box(slide1, Inches(1), Inches(1.5), Inches(11.3), Inches(1.2),
             "CNEAv5 Neural Interfacing Platform", font_size=44, color=CYAN, bold=True,
             alignment=PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide1, Inches(1), Inches(2.8), Inches(11.3), Inches(0.8),
             "Technical Architecture Overview", font_size=28, color=WHITE,
             alignment=PP_ALIGN.CENTER)

# Description
add_text_box(slide1, Inches(2), Inches(4.0), Inches(9.3), Inches(1.0),
             "Real-time Neural Data Acquisition, Visualization & Analysis Platform\n"
             "4096-Channel Electrode Array  |  FPGA Hardware Integration  |  AI-Powered Agents",
             font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Decorative line
line_shape = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(3.7), Inches(5.3), Pt(2))
line_shape.fill.solid()
line_shape.fill.fore_color.rgb = CYAN
line_shape.line.fill.background()

# Footer
add_text_box(slide1, Inches(1), Inches(6.5), Inches(11.3), Inches(0.5),
             "Eminence Tech Solutions  |  Confidential", font_size=12, color=MID_GRAY,
             alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SLIDE 2 — System Architecture Overview (High-Level)
# ════════════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide2, BG_DARK)

add_text_box(slide2, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             "System Architecture Overview", font_size=30, color=CYAN, bold=True)

# ── Layer boxes (top to bottom) ──

# Row Y positions
y_user    = Inches(1.2)
y_proxy   = Inches(2.2)
y_front   = Inches(3.0)
y_bff     = Inches(3.8)
y_backend = Inches(4.6)
y_agents  = Inches(5.5)
y_data    = Inches(6.3)

row_h = Inches(0.65)

layers = [
    ("Users / Browsers (Desktop & Mobile)", y_user, RGBColor(0x37, 0x5A, 0x7F)),
    ("Nginx SSL Reverse Proxy  [:3025 HTTPS / :80 HTTP Redirect]", y_proxy, RGBColor(0x4A, 0x3A, 0x2A)),
    ("React Frontend  [TypeScript + Vite + Redux Toolkit + TailwindCSS v4]", y_front, RGBColor(0x1E, 0x3A, 0x5F)),
    ("Node.js BFF  [Express + Socket.io — WebSocket Relay & REST Proxy — :3026]", y_bff, RGBColor(0x2A, 0x4A, 0x2A)),
    ("Django Backend  [Daphne/ASGI + DRF + Channels + Celery — :8085]", y_backend, RGBColor(0x2A, 0x3A, 0x2A)),
    ("Agent Cluster  [7 FastAPI Microservices + MCP Protocol — :8088-8094]", y_agents, RGBColor(0x3A, 0x2A, 0x4A)),
    ("Data Layer  [PostgreSQL 16 + TimescaleDB + pgvector  |  Redis 7]", y_data, RGBColor(0x4A, 0x2A, 0x1E)),
]

for label, y, color in layers:
    box = add_shape(slide2, Inches(1.0), y, Inches(11.3), row_h, color, BORDER_COLOR)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(4)

# Side labels
side_labels = [
    ("HTTPS", Inches(1.65)),
    ("Internal", Inches(2.55)),
    ("WebSocket\n+ REST", Inches(3.35)),
    ("REST API\n+ Channels", Inches(4.15)),
    ("MCP + Redis\nPub/Sub", Inches(5.0)),
    ("SQL + Redis", Inches(5.85)),
]
for text, y in side_labels:
    add_text_box(slide2, Inches(12.4), y, Inches(1.0), Inches(0.6),
                 text, font_size=8, color=MID_GRAY, alignment=PP_ALIGN.LEFT)

# External services box
ext = add_shape(slide2, Inches(0.15), Inches(5.5), Inches(0.7), Inches(1.4), BG_CARD, BORDER_COLOR)
tf = ext.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "FPGA\nHardware\n(Opal Kelly)"
p.font.size = Pt(8)
p.font.color.rgb = AMBER
p.font.bold = True
p.font.name = "Segoe UI"
p.alignment = PP_ALIGN.CENTER

# Langfuse
lang = add_shape(slide2, Inches(0.15), Inches(4.2), Inches(0.7), Inches(0.8), BG_CARD, BORDER_COLOR)
tf = lang.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Langfuse\nLLM Obs\n:8095"
p.font.size = Pt(8)
p.font.color.rgb = PURPLE
p.font.bold = True
p.font.name = "Segoe UI"
p.alignment = PP_ALIGN.CENTER

# Ollama
oll = add_shape(slide2, Inches(0.15), Inches(3.4), Inches(0.7), Inches(0.7), BG_CARD, BORDER_COLOR)
tf = oll.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Ollama\nLLM\nHost"
p.font.size = Pt(8)
p.font.color.rgb = GREEN
p.font.bold = True
p.font.name = "Segoe UI"
p.alignment = PP_ALIGN.CENTER


# ════════════════════════════════════════════════════════════════════
# SLIDE 3 — Detailed Service Architecture
# ════════════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide3, BG_DARK)

add_text_box(slide3, Inches(0.5), Inches(0.2), Inches(12), Inches(0.5),
             "Docker Compose Services (14 Containers)", font_size=28, color=CYAN, bold=True)

# Service cards
services = [
    # (name, port, description, color, x, y)
    ("nginx",         ":3025/:80",  "SSL Reverse Proxy\nTLS 1.2/1.3, HSTS",           AMBER,  0.3,  1.0),
    ("frontend",      "Internal",   "React + Vite + TS\nRedux, TailwindCSS v4",        CYAN,   2.6,  1.0),
    ("bff",           ":3026",      "Node.js Express\nSocket.io + Redis",               GREEN,  4.9,  1.0),
    ("django",        ":8085",      "Daphne/ASGI\nDRF + Channels + MCP",               GREEN,  7.2,  1.0),
    ("postgres",      ":5435",      "TimescaleDB PG16\npgvector embeddings",            RED,    9.5,  1.0),
    ("redis",         ":6385",      "Redis 7 Alpine\nPub/Sub + Cache",                  RED,    11.5, 1.0),

    ("data-acquisition", ":8088",   "FPGA Data Stream\nChannel Sampling",               PURPLE, 0.3,  3.0),
    ("signal-processing",":8089",   "Spike Detection\nFFT, Filtering",                  PURPLE, 2.6,  3.0),
    ("hardware-control", ":8090",   "FPGA Config\nBias/Clock/TIA",                      PURPLE, 4.9,  3.0),
    ("storage",       ":8091",      "Recording Mgmt\nHDF5/CSV Export",                  PURPLE, 7.2,  3.0),
    ("ai-ml",         ":8092",      "Spike Sorting\nPattern Recognition",               PURPLE, 9.5,  3.0),
    ("notification",  ":8093",      "Alerts & Email\nThreshold Monitor",                PURPLE, 11.5, 3.0),
    ("llm",           ":8094",      "LangGraph Bot\nRAG + Ollama",                      PURPLE, 0.3,  4.7),
    ("langfuse",      ":8095",      "LLM Observability\nTracing & Metrics",             PURPLE, 2.6,  4.7),
]

for name, port, desc, color, x, y in services:
    card = add_shape(slide3, Inches(x), Inches(y), Inches(2.1), Inches(1.6), BG_CARD, color, Pt(1.5))

    # Service name
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(6)
    tf.margin_left = Pt(8)
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(11)
    p.font.color.rgb = color
    p.font.bold = True
    p.font.name = "Consolas"

    p2 = tf.add_paragraph()
    p2.text = port
    p2.font.size = Pt(9)
    p2.font.color.rgb = MID_GRAY
    p2.font.name = "Consolas"
    p2.space_before = Pt(2)

    p3 = tf.add_paragraph()
    p3.text = desc
    p3.font.size = Pt(9)
    p3.font.color.rgb = LIGHT_GRAY
    p3.font.name = "Segoe UI"
    p3.space_before = Pt(6)

# Group label: Agent Cluster
add_text_box(slide3, Inches(0.3), Inches(2.65), Inches(3.0), Inches(0.35),
             "Agent Cluster (FastAPI + MCP)", font_size=12, color=PURPLE, bold=True)

# Group label: Infrastructure
add_text_box(slide3, Inches(0.3), Inches(0.7), Inches(3.0), Inches(0.3),
             "Infrastructure & Application Services", font_size=12, color=CYAN, bold=True)

# Connection info
add_text_box(slide3, Inches(5.0), Inches(4.8), Inches(8.0), Inches(1.5),
             "Data Flow: FPGA --> Data Acquisition Agent --> Redis Pub/Sub --> Signal Processing --> Django --> BFF (WebSocket) --> Frontend\n"
             "Control Flow: Frontend --> BFF --> Django REST API --> Agent MCP Tools --> Hardware\n"
             "Chat Flow: Frontend --> BFF (WebSocket) --> Django --> LLM Agent (LangGraph) --> Ollama/OpenAI/Anthropic\n"
             "All agents register tools via MCP (Model Context Protocol) with Django MCP Server",
             font_size=10, color=MID_GRAY)


# ════════════════════════════════════════════════════════════════════
# SLIDE 4 — Agent Architecture & MCP
# ════════════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide4, BG_DARK)

add_text_box(slide4, Inches(0.5), Inches(0.2), Inches(12), Inches(0.5),
             "Agent Architecture & MCP Protocol", font_size=28, color=CYAN, bold=True)

# MCP Server box (center)
mcp_box = add_shape(slide4, Inches(5.0), Inches(2.5), Inches(3.3), Inches(2.5), BG_SURFACE, CYAN, Pt(2))
tf = mcp_box.text_frame
tf.word_wrap = True
tf.margin_top = Pt(10)
p = tf.paragraphs[0]
p.text = "Django MCP Server"
p.font.size = Pt(16)
p.font.color.rgb = CYAN
p.font.bold = True
p.font.name = "Segoe UI"
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "\nTool Registry\nTool Execution\nAgent Discovery\nHealth Monitoring"
p2.font.size = Pt(11)
p2.font.color.rgb = LIGHT_GRAY
p2.font.name = "Segoe UI"
p2.alignment = PP_ALIGN.CENTER

# Agent boxes around the MCP server
agent_details = [
    ("Data Acquisition\n:8088", "start_recording\nstop_recording\nget_channel_data\nset_sample_rate", 0.3, 1.0),
    ("Signal Processing\n:8089", "detect_spikes\napply_filter\nrun_fft\nget_spike_stats", 0.3, 3.5),
    ("Hardware Control\n:8090", "set_bias_config\nset_clock_config\nset_tia_gain\nset_stim_params", 0.3, 5.8),
    ("Storage\n:8091", "save_recording\nexport_hdf5\nexport_csv\nlist_recordings", 9.3, 1.0),
    ("AI/ML\n:8092", "sort_spikes\nclassify_pattern\ndetect_anomaly\ntrain_model", 9.3, 3.5),
    ("Notification\n:8093", "send_alert\nset_threshold\nget_alert_history", 9.3, 5.8),
    ("LLM (LangGraph)\n:8094", "chat_query\nrag_search\nembed_document\nget_context", 5.0, 5.8),
]

for name, tools, x, y in agent_details:
    card = add_shape(slide4, Inches(x), Inches(y), Inches(3.5), Inches(1.6), BG_CARD, PURPLE, Pt(1.5))
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(6)
    tf.margin_left = Pt(8)
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(11)
    p.font.color.rgb = PURPLE
    p.font.bold = True
    p.font.name = "Consolas"

    p2 = tf.add_paragraph()
    p2.text = "MCP Tools:"
    p2.font.size = Pt(8)
    p2.font.color.rgb = CYAN
    p2.font.bold = True
    p2.font.name = "Segoe UI"
    p2.space_before = Pt(4)

    p3 = tf.add_paragraph()
    p3.text = tools
    p3.font.size = Pt(8)
    p3.font.color.rgb = MID_GRAY
    p3.font.name = "Consolas"
    p3.space_before = Pt(2)


# ════════════════════════════════════════════════════════════════════
# SLIDE 5 — Frontend Architecture
# ════════════════════════════════════════════════════════════════════
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide5, BG_DARK)

add_text_box(slide5, Inches(0.5), Inches(0.2), Inches(12), Inches(0.5),
             "Frontend Architecture", font_size=28, color=CYAN, bold=True)

# Pages section
pages_data = [
    ("Dashboard", "System overview, agent health,\nrecording stats, quick actions"),
    ("Visualization", "Real-time 4-panel view:\nHeatmap, Raster, Spectrum, Waveform"),
    ("Controls", "FPGA configuration presets,\nbias/clock/TIA/gain/stim params"),
    ("Recordings", "Browse, search, filter,\nexport recordings (HDF5/CSV)"),
    ("Experiments", "Experiment management,\nsession tracking, annotations"),
    ("Settings", "User profile, theme,\nnotification preferences"),
]

add_text_box(slide5, Inches(0.3), Inches(0.9), Inches(3), Inches(0.4),
             "Pages (React Router)", font_size=14, color=GREEN, bold=True)

for i, (name, desc) in enumerate(pages_data):
    x = Inches(0.3 + (i % 3) * 4.2)
    y = Inches(1.4 + (i // 3) * 1.6)
    card = add_shape(slide5, x, y, Inches(3.8), Inches(1.3), BG_CARD, GREEN, Pt(1))
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(8)
    tf.margin_left = Pt(10)
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(13)
    p.font.color.rgb = GREEN
    p.font.bold = True
    p.font.name = "Segoe UI"
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(10)
    p2.font.color.rgb = LIGHT_GRAY
    p2.font.name = "Segoe UI"
    p2.space_before = Pt(4)

# State Management section
add_text_box(slide5, Inches(0.3), Inches(4.6), Inches(6), Inches(0.4),
             "Redux Toolkit State Management", font_size=14, color=AMBER, bold=True)

slices = [
    ("recordingSlice", "isRecording, status, duration,\nspikeCount, dataRate, bufferUsage"),
    ("configSlice", "bias, clock, TIA, gain,\nstimulation, pixel, presets"),
    ("visualizationSlice", "viewType, activeChannels,\ncolorScale, timeWindow"),
    ("agentsSlice", "agents[], status, health,\nmetrics per agent"),
    ("chatSlice", "messages[], isPanelOpen,\nisLoading, context"),
]

for i, (name, desc) in enumerate(slices):
    x = Inches(0.3 + i * 2.5)
    card = add_shape(slide5, x, Inches(5.1), Inches(2.3), Inches(1.3), BG_CARD, AMBER, Pt(1))
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(6)
    tf.margin_left = Pt(6)
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(9)
    p.font.color.rgb = AMBER
    p.font.bold = True
    p.font.name = "Consolas"
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(8)
    p2.font.color.rgb = LIGHT_GRAY
    p2.font.name = "Segoe UI"
    p2.space_before = Pt(4)


# ════════════════════════════════════════════════════════════════════
# SLIDE 6 — Data Flow & Communication
# ════════════════════════════════════════════════════════════════════
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide6, BG_DARK)

add_text_box(slide6, Inches(0.5), Inches(0.2), Inches(12), Inches(0.5),
             "Data Flow & Communication Patterns", font_size=28, color=CYAN, bold=True)

# Flow 1: Real-time Data
add_text_box(slide6, Inches(0.3), Inches(1.0), Inches(12.5), Inches(0.4),
             "1. Real-Time Neural Data Flow", font_size=16, color=GREEN, bold=True)

flow1_steps = ["FPGA\nHardware", "Data Acq.\nAgent", "Redis\nPub/Sub", "Signal Proc.\nAgent", "Django\nChannels", "BFF\nSocket.io", "React\nFrontend"]
for i, step in enumerate(flow1_steps):
    x = Inches(0.3 + i * 1.8)
    card = add_shape(slide6, x, Inches(1.5), Inches(1.5), Inches(0.9), BG_CARD, GREEN, Pt(1))
    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = step
    p.font.size = Pt(9)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(6)

# Flow 2: Control Commands
add_text_box(slide6, Inches(0.3), Inches(2.7), Inches(12.5), Inches(0.4),
             "2. Hardware Control Flow", font_size=16, color=AMBER, bold=True)

flow2_steps = ["User\nInteraction", "React\nFrontend", "BFF\nREST Proxy", "Django\nREST API", "MCP Tool\nExecution", "Hardware\nAgent", "FPGA\nHardware"]
for i, step in enumerate(flow2_steps):
    x = Inches(0.3 + i * 1.8)
    card = add_shape(slide6, x, Inches(3.2), Inches(1.5), Inches(0.9), BG_CARD, AMBER, Pt(1))
    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = step
    p.font.size = Pt(9)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(6)

# Flow 3: Chat / LLM
add_text_box(slide6, Inches(0.3), Inches(4.4), Inches(12.5), Inches(0.4),
             "3. AI Chat Flow (LangGraph + RAG)", font_size=16, color=PURPLE, bold=True)

flow3_steps = ["User\nQuery", "BFF\nWebSocket", "Django\nRouting", "LLM Agent\nLangGraph", "RAG\npgvector", "Ollama /\nOpenAI", "Response\nStream"]
for i, step in enumerate(flow3_steps):
    x = Inches(0.3 + i * 1.8)
    card = add_shape(slide6, x, Inches(4.9), Inches(1.5), Inches(0.9), BG_CARD, PURPLE, Pt(1))
    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = step
    p.font.size = Pt(9)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(6)

# Protocol labels
protocols = [
    ("USB 3.0",        Inches(1.35), Inches(1.42), GREEN),
    ("Redis Pub/Sub",  Inches(4.95), Inches(1.42), GREEN),
    ("WebSocket",      Inches(9.75), Inches(1.42), GREEN),
    ("HTTPS",          Inches(1.35), Inches(3.12), AMBER),
    ("REST API",       Inches(4.95), Inches(3.12), AMBER),
    ("MCP",            Inches(8.55), Inches(3.12), AMBER),
    ("Socket.io",      Inches(1.35), Inches(4.82), PURPLE),
    ("pgvector SQL",   Inches(6.75), Inches(4.82), PURPLE),
    ("HTTP/gRPC",      Inches(9.75), Inches(4.82), PURPLE),
]

for label, x, y, color in protocols:
    add_text_box(slide6, x, y, Inches(1.5), Inches(0.3), label, font_size=7, color=color, bold=True,
                 alignment=PP_ALIGN.CENTER, font_name="Consolas")

# WebSocket Channels
add_text_box(slide6, Inches(0.3), Inches(6.1), Inches(12.5), Inches(0.5),
             "WebSocket Channels: recording:data | recording:spikes | chat:message | agents:status | visualization:stream",
             font_size=11, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SLIDE 7 — Technology Stack
# ════════════════════════════════════════════════════════════════════
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide7, BG_DARK)

add_text_box(slide7, Inches(0.5), Inches(0.2), Inches(12), Inches(0.5),
             "Technology Stack", font_size=28, color=CYAN, bold=True)

stack_categories = [
    ("Frontend", CYAN, [
        "React 18 + TypeScript",
        "Vite Build Tool",
        "Redux Toolkit",
        "TailwindCSS v4",
        "Recharts / D3.js",
        "Socket.io Client",
        "Lucide React Icons",
        "React Router v6",
    ]),
    ("BFF Layer", GREEN, [
        "Node.js 18+",
        "Express.js",
        "Socket.io Server",
        "ioredis Client",
        "Axios HTTP Client",
        "CORS Middleware",
    ]),
    ("Backend", GREEN, [
        "Python 3.11",
        "Django 5 + DRF",
        "Django Channels",
        "Daphne (ASGI)",
        "Celery Task Queue",
        "djangorestframework",
        "django-cors-headers",
    ]),
    ("AI / Agents", PURPLE, [
        "FastAPI",
        "LangGraph",
        "LangChain",
        "Ollama (Local LLM)",
        "OpenAI / Anthropic",
        "pgvector (RAG)",
        "Langfuse (Obs.)",
    ]),
    ("Data Layer", RED, [
        "PostgreSQL 16",
        "TimescaleDB",
        "pgvector Extension",
        "Redis 7",
        "HDF5 / CSV Export",
    ]),
    ("Infrastructure", AMBER, [
        "Docker Compose",
        "Nginx (SSL Proxy)",
        "TLS 1.2/1.3",
        "Let's Encrypt",
        "Ubuntu / Linux",
    ]),
]

for i, (cat, color, items) in enumerate(stack_categories):
    x = Inches(0.3 + (i % 6) * 2.15)
    y = Inches(1.0)

    card = add_shape(slide7, x, y, Inches(2.0), Inches(5.5), BG_CARD, color, Pt(1.5))
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(10)
    tf.margin_left = Pt(10)

    p = tf.paragraphs[0]
    p.text = cat
    p.font.size = Pt(13)
    p.font.color.rgb = color
    p.font.bold = True
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER

    for item in items:
        pi = tf.add_paragraph()
        pi.text = f"  {item}"
        pi.font.size = Pt(9)
        pi.font.color.rgb = LIGHT_GRAY
        pi.font.name = "Segoe UI"
        pi.space_before = Pt(6)


# ════════════════════════════════════════════════════════════════════
# SLIDE 8 — Database Schema
# ════════════════════════════════════════════════════════════════════
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide8, BG_DARK)

add_text_box(slide8, Inches(0.5), Inches(0.2), Inches(12), Inches(0.5),
             "Database Architecture", font_size=28, color=CYAN, bold=True)

tables = [
    ("recordings", "id, name, experiment_id,\nstart_time, end_time, status,\nsample_rate, channel_count,\nfile_path, file_size, metadata", 0.3, 1.0),
    ("experiments", "id, name, description,\nprotocol, status, created_by,\ncreated_at, updated_at,\nnotes, parameters", 4.6, 1.0),
    ("neural_data\n(TimescaleDB Hypertable)", "time, recording_id,\nchannel_id, voltage,\nfiltered_voltage,\nsampled at configurable rate", 8.9, 1.0),
    ("spike_events", "id, recording_id,\nchannel_id, timestamp,\namplitude, waveform,\ncluster_id, sorted_label", 0.3, 3.8),
    ("electrode_config", "id, recording_id,\nchannel_id, x_pos, y_pos,\nimpedance, enabled,\nbias_settings", 4.6, 3.8),
    ("users", "id, email, name,\nrole (Admin/Researcher/\nOperator/Viewer),\npassword_hash, tokens", 8.9, 3.8),
]

for name, cols, x, y in tables:
    card = add_shape(slide8, Inches(x), Inches(y), Inches(3.8), Inches(2.4), BG_CARD, RED, Pt(1.5))
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(8)
    tf.margin_left = Pt(10)
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(12)
    p.font.color.rgb = RED
    p.font.bold = True
    p.font.name = "Consolas"
    p2 = tf.add_paragraph()
    p2.text = cols
    p2.font.size = Pt(9)
    p2.font.color.rgb = LIGHT_GRAY
    p2.font.name = "Consolas"
    p2.space_before = Pt(8)

# Redis section
add_text_box(slide8, Inches(0.3), Inches(6.4), Inches(12.5), Inches(0.6),
             "Redis 7:  Pub/Sub Channels (recording_data, spike_events, agent_status)  |  "
             "Django Channel Layer  |  Celery Broker  |  Session Cache  |  Rate Limiting",
             font_size=11, color=AMBER, alignment=PP_ALIGN.CENTER)

# pgvector note
add_text_box(slide8, Inches(0.3), Inches(7.0), Inches(12.5), Inches(0.4),
             "pgvector Extension: 1536-dimensional embeddings for RAG pipeline (LLM Agent document search)",
             font_size=10, color=PURPLE, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SLIDE 9 — Deployment Architecture
# ════════════════════════════════════════════════════════════════════
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide9, BG_DARK)

add_text_box(slide9, Inches(0.5), Inches(0.2), Inches(12), Inches(0.5),
             "Deployment & Security", font_size=28, color=CYAN, bold=True)

# Network diagram
add_text_box(slide9, Inches(0.3), Inches(0.9), Inches(6), Inches(0.4),
             "Network & SSL Configuration", font_size=16, color=GREEN, bold=True)

net_items = [
    ("Internet / Browser", "HTTPS :3025", 0.3, 1.5, LIGHT_GRAY),
    ("Nginx SSL Proxy", "TLS 1.2/1.3\nHSTS, OCSP Stapling\nModern Cipher Suite", 3.5, 1.5, AMBER),
    ("Frontend\n(Internal :80)", "Static React SPA\nServed by Nginx", 6.7, 1.5, CYAN),
    ("BFF\n(Internal :3026)", "Express + Socket.io\nWebSocket Upgrade", 3.5, 3.5, GREEN),
    ("Django\n(Internal :8085)", "Daphne ASGI\n/admin/ + /mcp/", 6.7, 3.5, GREEN),
]

for name, desc, x, y, color in net_items:
    card = add_shape(slide9, Inches(x), Inches(y), Inches(2.8), Inches(1.6), BG_CARD, color, Pt(1.5))
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(6)
    tf.margin_left = Pt(8)
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(11)
    p.font.color.rgb = color
    p.font.bold = True
    p.font.name = "Segoe UI"
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(9)
    p2.font.color.rgb = MID_GRAY
    p2.font.name = "Segoe UI"
    p2.space_before = Pt(4)

# Security section
add_text_box(slide9, Inches(0.3), Inches(5.4), Inches(6), Inches(0.4),
             "Security Features", font_size=16, color=RED, bold=True)

security_items = [
    "JWT Authentication with Refresh Tokens",
    "Role-Based Access Control (Admin, Researcher, Operator, Viewer)",
    "SSL/TLS 1.2/1.3 Encryption with HSTS",
    "CORS Configuration (whitelisted origins)",
    "Django ORM (SQL Injection Prevention)",
    "Input Validation & Sanitization",
    "HTTP Security Headers (X-Frame-Options, X-Content-Type, CSP)",
    "Redis Password Authentication",
    "Docker Network Isolation",
]

for i, item in enumerate(security_items):
    x = Inches(0.5 + (i // 5) * 6.3)
    y = Inches(5.9 + (i % 5) * 0.28)
    add_text_box(slide9, x, y, Inches(6.0), Inches(0.28),
                 f"  {item}", font_size=9, color=LIGHT_GRAY)

# Port mapping table area
add_text_box(slide9, Inches(9.7), Inches(0.9), Inches(3.3), Inches(0.4),
             "Port Mapping", font_size=14, color=AMBER, bold=True)

ports = [
    ("3025", "HTTPS (Nginx)"),
    ("80", "HTTP Redirect"),
    ("3026", "BFF (Internal)"),
    ("8085", "Django (Internal)"),
    ("5435", "PostgreSQL"),
    ("6385", "Redis"),
    ("8088-8094", "Agent Cluster"),
    ("8095", "Langfuse"),
]

for i, (port, desc) in enumerate(ports):
    y = Inches(1.4 + i * 0.35)
    add_text_box(slide9, Inches(9.7), y, Inches(1.0), Inches(0.3),
                 port, font_size=9, color=CYAN, bold=True, font_name="Consolas")
    add_text_box(slide9, Inches(10.8), y, Inches(2.2), Inches(0.3),
                 desc, font_size=9, color=LIGHT_GRAY)


# ════════════════════════════════════════════════════════════════════
# SLIDE 10 — Summary
# ════════════════════════════════════════════════════════════════════
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide10, BG_DARK)

add_text_box(slide10, Inches(1), Inches(0.5), Inches(11.3), Inches(0.8),
             "Architecture Summary", font_size=36, color=CYAN, bold=True,
             alignment=PP_ALIGN.CENTER)

summary_items = [
    ("14 Docker Services", "Fully containerized microservices architecture with Docker Compose orchestration"),
    ("7 AI-Powered Agents", "FastAPI microservices with MCP protocol for tool registration and execution"),
    ("4096 Channels", "64x64 electrode array with FPGA hardware interface (Opal Kelly USB 3.0)"),
    ("Real-Time Streaming", "Sub-second data flow via Redis Pub/Sub and WebSocket (Socket.io)"),
    ("LLM-Powered Assistant", "LangGraph chatbot with RAG pipeline, supporting Ollama/OpenAI/Anthropic"),
    ("4 Visualization Types", "Heatmap, Raster Plot, Frequency Spectrum, Waveform — all real-time"),
    ("Role-Based Security", "JWT auth, 4 user roles, SSL/TLS, CORS, input validation"),
    ("Mobile Responsive", "Full mobile support with adaptive layouts and touch-friendly controls"),
]

for i, (title, desc) in enumerate(summary_items):
    x = Inches(0.5 + (i % 2) * 6.4)
    y = Inches(1.5 + (i // 2) * 1.4)
    card = add_shape(slide10, x, y, Inches(6.0), Inches(1.15), BG_CARD, BORDER_COLOR)
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(8)
    tf.margin_left = Pt(12)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.color.rgb = CYAN
    p.font.bold = True
    p.font.name = "Segoe UI"
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(11)
    p2.font.color.rgb = LIGHT_GRAY
    p2.font.name = "Segoe UI"
    p2.space_before = Pt(4)

# Footer
add_text_box(slide10, Inches(1), Inches(6.8), Inches(11.3), Inches(0.5),
             "CNEAv5 Neural Interfacing Platform  |  Eminence Tech Solutions",
             font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# ── Save ────────────────────────────────────────────────────────────
output_path = "/home/user/capstone-project1/docs/CNEAv5_Technical_Architecture.pptx"
prs.save(output_path)
print(f"Presentation saved to {output_path}")
print(f"Slides: {len(prs.slides)}")
